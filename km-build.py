#!/usr/bin/env python3
"""
km-build: 掃描指定目錄，用 Claude API 為每個文件產生 .md 摘要，
建立 index.md 和反向連結，產生 graph-data.json。

用法：
    python3 km-build.py /path/to/your/folder
    python3 km-build.py /path/to/your/folder --skip-llm   # 只重建索引，不呼叫 API
"""

import argparse
import base64
import hashlib
import json
import os
import re
import sys
import time
from datetime import datetime
from pathlib import Path

from dotenv import load_dotenv

# 載入 .env
load_dotenv(Path(__file__).parent / '.env')

# ====== 設定 ======

# 需要 LLM 處理的檔案類型
LLM_EXTENSIONS = {'.pdf', '.doc', '.docx', '.ppt', '.pptx', '.txt', '.rtf', '.md'}

# 所有支援的檔案類型分類
FILE_CATEGORIES = {
    '文件': {'.pdf', '.doc', '.docx', '.txt', '.rtf', '.odt', '.md'},
    '簡報': {'.ppt', '.pptx', '.key', '.odp'},
    '試算表': {'.xls', '.xlsx', '.csv', '.ods', '.numbers'},
    '圖片': {'.jpg', '.jpeg', '.png', '.gif', '.svg', '.bmp', '.webp', '.tiff'},
    '影音': {'.mp4', '.mov', '.avi', '.mkv', '.mp3', '.wav', '.flac', '.m4a'},
    '程式碼': {'.py', '.js', '.ts', '.html', '.css', '.json', '.xml', '.yaml', '.yml'},
    '壓縮檔': {'.zip', '.rar', '.7z', '.tar', '.gz', '.bz2'},
}

EXT_TO_CATEGORY = {}
for cat, exts in FILE_CATEGORIES.items():
    for ext in exts:
        EXT_TO_CATEGORY[ext] = cat

# Claude API 支援的文件 MIME 類型
MIME_TYPES = {
    '.pdf': 'application/pdf',
    '.txt': 'text/plain',
    '.md': 'text/plain',
    '.rtf': 'text/plain',
    '.csv': 'text/csv',
}

# 需要轉成文字再傳的類型（Claude 不直接支援）
TEXT_READABLE = {'.txt', '.md', '.rtf', '.csv'}

# PDF 大小上限（30MB，Claude API 限制為 32MB）
MAX_PDF_SIZE = 30 * 1024 * 1024
# 文字檔大小上限
MAX_TEXT_SIZE = 10 * 1024 * 1024


def get_client():
    """取得 Anthropic client"""
    import anthropic
    api_key = os.getenv('ANTHROPIC_API_KEY')
    if not api_key:
        print('❌ 請在 .env 檔設定 ANTHROPIC_API_KEY')
        sys.exit(1)
    return anthropic.Anthropic(api_key=api_key)


def stable_id(path: str) -> str:
    return hashlib.md5(path.encode('utf-8')).hexdigest()[:10]


def format_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes} B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f} KB"
    elif size_bytes < 1024 * 1024 * 1024:
        return f"{size_bytes / (1024 * 1024):.1f} MB"
    else:
        return f"{size_bytes / (1024 * 1024 * 1024):.1f} GB"


def get_file_category(ext: str) -> str:
    return EXT_TO_CATEGORY.get(ext.lower(), '其他')


def extract_docx_text(filepath: Path) -> str:
    """從 .docx 提取文字"""
    from docx import Document
    doc = Document(filepath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # 也提取表格
    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                paragraphs.append(' | '.join(cells))
    return '\n\n'.join(paragraphs)


def extract_pptx_text(filepath: Path) -> str:
    """從 .pptx 提取文字"""
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f'--- Slide {i} ---\n' + '\n'.join(texts))
    return '\n\n'.join(slides)


def extract_xlsx_text(filepath: Path) -> str:
    """從 .xlsx 提取文字"""
    from openpyxl import load_workbook
    wb = load_workbook(filepath, data_only=True)
    sheets = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(max_row=200, values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                rows.append(' | '.join(cells))
        if rows:
            sheets.append(f'--- Sheet: {ws.title} ---\n' + '\n'.join(rows))
    return '\n\n'.join(sheets)


def read_file_content(filepath: Path) -> tuple:
    """
    讀取檔案內容，回傳 (content_type, content)。
    content_type: 'text', 'pdf', 'skip'
    """
    ext = filepath.suffix.lower()
    size = filepath.stat().st_size

    if size == 0:
        return 'skip', '空檔案'

    # 文字檔
    if ext in TEXT_READABLE:
        if size > MAX_TEXT_SIZE:
            return 'skip', f'檔案太大（{format_size(size)}）'
        try:
            text = filepath.read_text(encoding='utf-8')
            return 'text', text
        except UnicodeDecodeError:
            try:
                text = filepath.read_text(encoding='big5')
                return 'text', text
            except Exception:
                return 'skip', '無法讀取文字編碼'

    # PDF
    if ext == '.pdf':
        if size > MAX_PDF_SIZE:
            return 'skip', f'檔案太大（{format_size(size)}）'
        data = filepath.read_bytes()
        b64 = base64.standard_b64encode(data).decode('utf-8')
        return 'pdf', b64

    # DOCX
    if ext in ('.doc', '.docx'):
        try:
            text = extract_docx_text(filepath)
            if text.strip():
                return 'text', text
            return 'skip', 'DOCX 無文字內容'
        except Exception as e:
            return 'skip', f'DOCX 解析失敗：{e}'

    # PPTX
    if ext in ('.ppt', '.pptx'):
        try:
            text = extract_pptx_text(filepath)
            if text.strip():
                return 'text', text
            return 'skip', 'PPTX 無文字內容'
        except Exception as e:
            return 'skip', f'PPTX 解析失敗：{e}'

    # XLSX
    if ext in ('.xls', '.xlsx'):
        try:
            text = extract_xlsx_text(filepath)
            if text.strip():
                return 'text', text
            return 'skip', 'XLSX 無文字內容'
        except Exception as e:
            return 'skip', f'XLSX 解析失敗：{e}'

    return 'skip', f'不支援的格式（{ext}）'


SYSTEM_PROMPT = """你是知識庫管理員。請分析提供的文件，用繁體中文回應。
嚴格按照以下 JSON 格式回應，不要加任何其他文字：

{
  "summary": "200-400字的摘要，涵蓋文件的主要內容、關鍵論點和結論",
  "keywords": ["關鍵字1", "關鍵字2", ...],
  "tags": ["#標籤1", "#標籤2", ...],
  "related_files": ["可能相關的檔案名稱1", "可能相關的檔案名稱2", ...]
}

related_files 請從提供的檔案清單中，根據主題相關性選出可能相關的檔案。"""

CLAUDE_MODELS = {
    'claude-sonnet': 'claude-sonnet-4-20250514',
    'claude-haiku': 'claude-haiku-4-5-20251001',
    'claude-opus': 'claude-opus-4-20250514',
}

GEMINI_MODELS = {
    'gemini-flash': 'gemini-2.5-flash',
    'gemini-pro': 'gemini-2.5-pro',
    'gemini-3-pro': 'gemini-3-pro-preview',
    'gemini-3-flash': 'gemini-3-flash-preview',
}


def call_llm(model_id: str, filepath: Path, all_filenames: list) -> dict:
    """呼叫 LLM 產生摘要。支援 Claude / Gemini / Ollama。"""
    content_type, content = read_file_content(filepath)

    if content_type == 'skip':
        return {
            'summary': f'（未解析：{content}）',
            'keywords': [],
            'tags': [],
            'related_files': [],
        }

    other_files = [f for f in all_filenames if f != filepath.name]
    files_list = '\n'.join(f'- {f}' for f in other_files[:50])
    user_prompt = f"檔案名稱：{filepath.name}\n\n同目錄下的其他檔案：\n{files_list}\n\n請分析此文件並回傳 JSON。"

    try:
        if model_id.startswith('claude'):
            text = _call_claude(model_id, content_type, content, user_prompt)
        elif model_id.startswith('gemini'):
            text = _call_gemini(model_id, content_type, content, user_prompt)
        elif model_id.startswith('ollama:'):
            text = _call_ollama(model_id.split(':', 1)[1], content_type, content, user_prompt)
        else:
            text = _call_claude('claude-sonnet', content_type, content, user_prompt)

        json_match = re.search(r'\{[\s\S]*\}', text)
        if json_match:
            return json.loads(json_match.group())
        else:
            return {'summary': text, 'keywords': [], 'tags': [], 'related_files': []}

    except Exception as e:
        print(f'  ⚠️ API 錯誤：{e}')
        return {
            'summary': f'（API 錯誤：{str(e)[:100]}）',
            'keywords': [],
            'tags': [],
            'related_files': [],
        }


def _call_claude(model_id: str, content_type: str, content: str, user_prompt: str) -> str:
    import anthropic
    client = anthropic.Anthropic(api_key=os.getenv('ANTHROPIC_API_KEY'))
    model_name = CLAUDE_MODELS.get(model_id, CLAUDE_MODELS['claude-sonnet'])

    user_content = []
    if content_type == 'pdf':
        user_content.append({"type": "document", "source": {"type": "base64", "media_type": "application/pdf", "data": content}})
    elif content_type == 'text':
        if len(content) > 50000:
            content = content[:50000] + '\n\n...（截斷）'
        user_content.append({"type": "text", "text": f"文件內容：\n\n{content}"})
    user_content.append({"type": "text", "text": user_prompt})

    response = client.messages.create(model=model_name, max_tokens=1024, system=SYSTEM_PROMPT, messages=[{"role": "user", "content": user_content}])
    return response.content[0].text.strip()


def _call_gemini(model_id: str, content_type: str, content: str, user_prompt: str) -> str:
    import urllib.request
    api_key = os.getenv('GOOGLE_API_KEY')
    model_name = GEMINI_MODELS.get(model_id, GEMINI_MODELS['gemini-flash'])

    text_content = content if content_type == 'text' else '（PDF 文件，僅提供檔名分析）'
    if len(text_content) > 50000:
        text_content = text_content[:50000] + '\n\n...（截斷）'

    payload = json.dumps({
        "contents": [{"parts": [{"text": f"{SYSTEM_PROMPT}\n\n文件內容：\n\n{text_content}\n\n{user_prompt}"}]}],
        "generationConfig": {"maxOutputTokens": 1024},
    }).encode('utf-8')

    req = urllib.request.Request(
        f'https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}',
        data=payload, headers={'Content-Type': 'application/json'}, method='POST')
    with urllib.request.urlopen(req, timeout=120) as resp:
        data = json.loads(resp.read())
        return data['candidates'][0]['content']['parts'][0]['text']


def _call_ollama(model: str, content_type: str, content: str, user_prompt: str) -> str:
    import urllib.request
    base = os.getenv('OLLAMA_BASE_URL', 'http://localhost:11434')

    text_content = content if content_type == 'text' else '（PDF 文件，僅提供檔名分析）'
    if len(text_content) > 30000:
        text_content = text_content[:30000] + '\n\n...（截斷）'

    payload = json.dumps({
        "model": model, "stream": False,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": f"文件內容：\n\n{text_content}\n\n{user_prompt}"},
        ],
    }).encode('utf-8')

    req = urllib.request.Request(f'{base}/api/chat', data=payload, headers={'Content-Type': 'application/json'}, method='POST')
    with urllib.request.urlopen(req, timeout=300) as resp:
        data = json.loads(resp.read())
        return data.get('message', {}).get('content', '')


def generate_md(filepath: Path, llm_result: dict, wiki_dir: Path) -> str:
    """為單一檔案產生 .md 內容"""
    stat = filepath.stat()
    ext = filepath.suffix.lower()
    category = get_file_category(ext)
    modified = datetime.fromtimestamp(stat.st_mtime).strftime('%Y-%m-%d')

    # YAML frontmatter
    tags_str = ', '.join(f'"{t}"' for t in llm_result.get('tags', []))
    keywords_str = ', '.join(f'"{k}"' for k in llm_result.get('keywords', []))

    md = f"""---
title: "{filepath.stem}"
source: "{filepath.name}"
source_path: "{filepath}"
category: "{category}"
ext: "{ext}"
size: "{format_size(stat.st_size)}"
modified: "{modified}"
tags: [{tags_str}]
keywords: [{keywords_str}]
---

# {filepath.stem}

| 欄位 | 值 |
|------|-----|
| 原始檔案 | `{filepath.name}` |
| 類型 | {category} ({ext}) |
| 大小 | {format_size(stat.st_size)} |
| 修改日期 | {modified} |

## 摘要

{llm_result.get('summary', '（無摘要）')}

## 關鍵字

{', '.join(llm_result.get('keywords', [])) or '（無）'}

## 相關文件

"""
    related = llm_result.get('related_files', [])
    if related:
        for rf in related:
            # 轉成 wiki link 格式（去掉副檔名）
            link_name = Path(rf).stem
            md += f'- [[{link_name}]]\n'
    else:
        md += '（無）\n'

    return md


def build_backlinks(wiki_dir: Path):
    """掃描所有 .md 的 [[連結]]，建立反向連結"""
    md_files = list(wiki_dir.glob('*.md'))
    if not md_files:
        return

    # 讀取所有 .md
    contents = {}
    for f in md_files:
        if f.name in ('index.md', 'log.md'):
            continue
        contents[f.name] = f.read_text(encoding='utf-8')

    # 掃描正向連結
    forward = {}
    for fname, content in contents.items():
        links = re.findall(r'\[\[([^\]]+)\]\]', content)
        forward[fname] = links

    # 建立反向連結表
    backlinks = {fname: [] for fname in contents}
    for src, links in forward.items():
        src_stem = Path(src).stem
        for link in links:
            target = link + '.md'
            if target in contents and target != src:
                if src_stem not in backlinks[target]:
                    backlinks[target].append(src_stem)

    # 寫入反向連結
    section_marker = '## 反向連結'
    for fname, content in contents.items():
        bl = backlinks[fname]

        # 移除舊的反向連結區塊
        content = re.sub(r'\n---\n## 反向連結.*', '', content, flags=re.DOTALL).rstrip()

        # 加入新的
        if bl:
            bl_text = '\n'.join(f'- [[{name}]]' for name in bl)
            content += f'\n\n---\n{section_marker}\n\n{bl_text}\n'
        else:
            content += f'\n\n---\n{section_marker}\n\n（無）\n'

        (wiki_dir / fname).write_text(content, encoding='utf-8')


def build_index(wiki_dir: Path, root_dir: Path):
    """建立 index.md — 所有頁面的目錄"""
    md_files = sorted(wiki_dir.glob('*.md'))

    lines = [f'# {root_dir.name} — 知識庫索引\n']
    lines.append(f'產生時間：{datetime.now().strftime("%Y-%m-%d %H:%M")}\n')
    lines.append(f'## 文件列表\n')
    lines.append('| 文件 | 類型 | 大小 | 修改日期 |')
    lines.append('|------|------|------|----------|')

    for f in md_files:
        if f.name in ('index.md', 'log.md'):
            continue
        # 從 frontmatter 讀取 metadata
        content = f.read_text(encoding='utf-8')
        category = _extract_fm(content, 'category') or ''
        size = _extract_fm(content, 'size') or ''
        modified = _extract_fm(content, 'modified') or ''
        stem = f.stem
        lines.append(f'| [[{stem}]] | {category} | {size} | {modified} |')

    (wiki_dir / 'index.md').write_text('\n'.join(lines) + '\n', encoding='utf-8')


def _extract_fm(content: str, key: str) -> str:
    """從 frontmatter 提取值"""
    match = re.search(rf'^{key}:\s*"?([^"\n]*)"?', content, re.MULTILINE)
    return match.group(1) if match else ''


def build_graph_data(wiki_dir: Path, root_dir: Path) -> dict:
    """從 .md 建立 graph-data.json"""
    nodes = []
    edges = []
    md_files = list(wiki_dir.glob('*.md'))

    # 為每個 .md 建立節點
    for f in md_files:
        if f.name in ('index.md', 'log.md'):
            continue

        content = f.read_text(encoding='utf-8')
        category = _extract_fm(content, 'category') or '其他'
        ext = _extract_fm(content, 'ext') or ''
        size = _extract_fm(content, 'size') or ''
        modified = _extract_fm(content, 'modified') or ''
        source_path = _extract_fm(content, 'source_path') or ''
        source = _extract_fm(content, 'source') or f.stem
        tags_match = re.search(r'^tags:\s*\[(.+)\]', content, re.MULTILINE)
        tags = re.findall(r'"([^"]+)"', tags_match.group(1)) if tags_match else []
        keywords_match = re.search(r'^keywords:\s*\[(.+)\]', content, re.MULTILINE)
        keywords = re.findall(r'"([^"]+)"', keywords_match.group(1)) if keywords_match else []

        # 提取摘要（## 摘要 之後的內容）
        summary_match = re.search(r'## 摘要\s*\n\s*(.+?)(?=\n## |\n---|\Z)', content, re.DOTALL)
        summary = summary_match.group(1).strip() if summary_match else ''

        node_id = stable_id(f.stem)
        nodes.append({
            'id': node_id,
            'label': source,
            'stem': f.stem,
            'type': 'file',
            'category': category,
            'ext': ext,
            'path': source_path,
            'size': size,
            'modified': modified,
            'modifiedDate': modified,
            'modifiedMonth': modified[:7] if modified else '',
            'tags': tags,
            'keywords': keywords,
            'summary': summary,
            'wikiFile': f.name,
        })

    # 從 [[連結]] 建立 edges
    stem_to_id = {n['stem']: n['id'] for n in nodes}
    for f in md_files:
        if f.name in ('index.md', 'log.md'):
            continue
        content = f.read_text(encoding='utf-8')
        links = re.findall(r'\[\[([^\]]+)\]\]', content)
        src_id = stem_to_id.get(f.stem)
        if not src_id:
            continue
        for link in links:
            tgt_id = stem_to_id.get(link)
            if tgt_id and tgt_id != src_id:
                # 避免重複邊
                edge_key = tuple(sorted([src_id, tgt_id]))
                if not any(tuple(sorted([e['source'], e['target']])) == edge_key for e in edges):
                    edges.append({'source': src_id, 'target': tgt_id, 'type': 'semantic'})

    # 統計
    categories = {}
    for n in nodes:
        cat = n['category']
        categories[cat] = categories.get(cat, 0) + 1

    dates = [n['modified'] for n in nodes if n['modified']]
    date_range = {'min': min(dates), 'max': max(dates)} if dates else None

    return {
        'meta': {
            'rootDir': str(root_dir),
            'rootName': root_dir.name,
            'scannedAt': datetime.now().isoformat(),
            'totalFiles': len(nodes),
            'totalFolders': 0,
            'categories': categories,
            'dateRange': date_range,
            'hasWiki': True,
        },
        'nodes': nodes,
        'edges': edges,
    }


def main():
    parser = argparse.ArgumentParser(description='掃描目錄，用 LLM 產生知識庫 .md')
    parser.add_argument('directory', help='要掃描的目錄路徑')
    parser.add_argument('-o', '--output', default='.', help='輸出目錄（預設為當前目錄）')
    parser.add_argument('--skip-llm', action='store_true', help='跳過 LLM，只重建索引和圖譜')
    parser.add_argument('--model', default='claude-sonnet', help='LLM 模型（claude-sonnet, claude-haiku, gemini-flash, gemini-pro, ollama:模型名）')
    args = parser.parse_args()

    root = Path(args.directory).resolve()
    if not root.is_dir():
        print(f'❌ 目錄不存在：{root}')
        sys.exit(1)

    output_dir = Path(args.output).resolve()
    # 用目錄名稱作為專案子目錄
    project_name = root.name
    wiki_dir = output_dir / 'wiki' / project_name
    wiki_dir.mkdir(parents=True, exist_ok=True)

    # 收集所有檔案
    all_files = []
    for dirpath, dirnames, filenames in os.walk(root):
        dirnames[:] = [d for d in dirnames if not d.startswith('.') and d not in {'node_modules', '__pycache__', '.git', 'venv'}]
        for fname in sorted(filenames):
            if fname.startswith('.') or fname.startswith('~$'):
                continue
            fpath = Path(dirpath) / fname
            all_files.append(fpath)

    all_filenames = [f.name for f in all_files]
    print(f'🔍 掃描目錄：{root}')
    print(f'   找到 {len(all_files)} 個檔案')

    if not args.skip_llm:
        # 篩選需要 LLM 處理的檔案
        llm_files = [f for f in all_files if f.suffix.lower() in LLM_EXTENSIONS]
        skip_files = [f for f in all_files if f.suffix.lower() not in LLM_EXTENSIONS]

        model_id = args.model
        print(f'   🤖 需要 LLM 處理：{len(llm_files)} 個（模型：{model_id}）')
        print(f'   ⏭️  跳過（圖片等）：{len(skip_files)} 個')

        # 為 LLM 檔案產生 .md
        for i, fpath in enumerate(llm_files, 1):
            md_path = wiki_dir / (fpath.stem + '.md')

            # 如果 .md 已存在且原檔沒更新，跳過
            if md_path.exists():
                md_mtime = md_path.stat().st_mtime
                src_mtime = fpath.stat().st_mtime
                if src_mtime <= md_mtime:
                    print(f'  ⏭️  [{i}/{len(llm_files)}] {fpath.name}（已存在，跳過）')
                    continue

            print(f'  🤖 [{i}/{len(llm_files)}] {fpath.name}...')
            llm_result = call_llm(model_id, fpath, all_filenames)
            md_content = generate_md(fpath, llm_result, wiki_dir)
            md_path.write_text(md_content, encoding='utf-8')
            print(f'     ✅ → {md_path.name}')

            # 避免 rate limit
            time.sleep(0.5)

        # 為非 LLM 檔案產生簡單 .md（只有 metadata，沒有摘要）
        for fpath in skip_files:
            md_path = wiki_dir / (fpath.stem + '.md')
            if md_path.exists():
                continue
            simple_result = {
                'summary': f'{get_file_category(fpath.suffix.lower())}檔案。',
                'keywords': [],
                'tags': [f'#{get_file_category(fpath.suffix.lower())}'],
                'related_files': [],
            }
            md_content = generate_md(fpath, simple_result, wiki_dir)
            md_path.write_text(md_content, encoding='utf-8')

    # 建立反向連結
    print('🔗 建立反向連結...')
    build_backlinks(wiki_dir)

    # 建立 index
    print('📇 建立 index.md...')
    build_index(wiki_dir, root)

    # 產生 graph-data.json
    print('📊 產生 graph-data.json...')
    graph_data = build_graph_data(wiki_dir, root)

    # graph-data.json 放在專案子目錄
    graph_file = wiki_dir / 'graph-data.json'
    with open(graph_file, 'w', encoding='utf-8') as f:
        json.dump(graph_data, f, ensure_ascii=False, indent=2)

    # 根目錄的 graph-data.json 由 switch_project 負責，這裡不寫

    print(f'\n✅ 完成！')
    print(f'   📁 專案：{project_name}')
    print(f'   📝 Wiki 頁面：{len(list(wiki_dir.glob("*.md")))} 個')
    print(f'   🔗 語意連結：{len(graph_data["edges"])} 條')
    print(f'   💾 {graph_file}')
    print(f'\n啟動方式：')
    print(f'   python3 km-server.py')
    print(f'   然後開啟 http://localhost:8765/graph-view.html')


if __name__ == '__main__':
    main()
